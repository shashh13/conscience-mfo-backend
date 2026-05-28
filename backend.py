"""
Conscience MFO — PPTX Report Generator Backend v2
Layout matches Investment_Report-Updated.pptx reference exactly.
"""

from flask import Flask, request, jsonify
import base64, io, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

app = Flask(__name__)

# ── Palette ───────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x2A, 0x4A)
GOLD       = RGBColor(0xC9, 0xA0, 0x2E)   # warm gold matching reference
GOLD_BG    = RGBColor(0xD4, 0xB4, 0x83)     # light gold banner fill
GOLD_HDR   = RGBColor(0xC8, 0xA9, 0x6E)   # header gold
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF9, 0xF9, 0xF9)
DARK_GREY  = RGBColor(0x2D, 0x2D, 0x2D)
MID_GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
BORDER     = RGBColor(0xD0, 0xD0, 0xD0)
GOLD_FILL  = RGBColor(0xD4, 0xAF, 0x70)

SECTOR_PAL = {
    "Technology":  RGBColor(0x1F, 0x4E, 0x79),
    "ETF":         RGBColor(0xC5, 0x7A, 0x1E),
    "Industrials": RGBColor(0x70, 0xAD, 0x47),
    "Financials":  RGBColor(0xFF, 0xC0, 0x00),
    "Healthcare":  RGBColor(0xFF, 0x00, 0x00),
    "Energy":      RGBColor(0x9D, 0xC3, 0xE6),
    "Real Estate": RGBColor(0x59, 0x59, 0x59),
}

# ── Helpers ───────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill, line_color=None, line_pt=0):
    from pptx.util import Pt as Pt2
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color and line_pt > 0:
        s.line.color.rgb = line_color; s.line.width = Pt2(line_pt)
    else:
        s.line.fill.background()
    return s

def add_text(slide, x, y, w, h, text, size=10, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_banner(slide, x, y, w, h, title, font_size=10):
    """Gold banner with bold centred title — matching reference style."""
    add_rect(slide, x, y, w, h, GOLD_FILL)
    add_text(slide, x, y, w, h, title,
             size=font_size, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)

def add_bullets(slide, x, y, w, h, bullets, size=9, color=DARK_GREY):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(4)
        run = p.add_run(); run.text = "•   " + b
        run.font.name = "Calibri"; run.font.size = Pt(size)
        run.font.color.rgb = color

# ── SLIDE 1 ───────────────────────────────────────────────────────────────────
def build_slide1(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = 13.33, 7.5
    allocs = sorted(data["allocations"], key=lambda x: x["alloc"], reverse=True)
    grand  = sum(a["alloc"] for a in allocs)

    # White background
    add_rect(slide, 0, 0, W, H, WHITE)

    # Title
    add_text(slide, 0.35, 0.18, 6, 0.55, "Allocation Plan",
             size=32, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT)

    # Gold banner
    add_rect(slide, 0.3, 0.88, W-0.6, 0.38, GOLD_FILL)
    add_text(slide, 0.3, 0.88, W-0.6, 0.38,
             "RECOMMENDED PORTFOLIO ALLOCATION",
             size=12, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)

    # Table geometry
    tbl_top  = 1.38
    row_h    = 0.268
    # col widths: ticker, company, alloc, wt% — two sets side by side
    cw = [0.75, 2.85, 0.72, 0.62]
    lx = 0.3
    rx = lx + sum(cw) + 0.28

    def draw_col_header(x):
        cx = x
        for i, lbl in enumerate(["TICKER", "COMPANY", "ALLOC", "WT%"]):
            add_rect(slide, cx, tbl_top, cw[i], row_h, GOLD_FILL)
            add_text(slide, cx+0.04, tbl_top+0.04, cw[i]-0.06, row_h-0.06,
                     lbl, size=9, bold=True, color=DARK_GREY, align=PP_ALIGN.CENTER)
            cx += cw[i]

    def draw_rows(items, sx):
        for idx, h in enumerate(items):
            ry  = tbl_top + row_h + idx * row_h
            bg  = WHITE if idx % 2 == 0 else LIGHT_GREY
            cx  = sx
            vals = [
                (h["ticker"],                            PP_ALIGN.CENTER, True),
                (h["name"],                              PP_ALIGN.CENTER, False),
                ("${:,}".format(int(h["alloc"])),        PP_ALIGN.CENTER, False),
                ("{:.1f}%".format(h["alloc"]/grand*100), PP_ALIGN.CENTER, False),
            ]
            for i, (val, al, bld) in enumerate(vals):
                add_rect(slide, cx, ry, cw[i], row_h, bg, BORDER, 0.3)
                add_text(slide, cx+0.04, ry+0.03, cw[i]-0.06, row_h-0.04,
                         val, size=9, bold=bld, color=DARK_GREY, align=al)
                cx += cw[i]

    draw_col_header(lx)
    draw_col_header(rx)
    draw_rows(allocs[:14], lx)
    draw_rows(allocs[14:], rx)

    # Sector key
    key_y = tbl_top + row_h * 15 + 0.12
    add_text(slide, lx, key_y, 1.0, 0.22, "SECTOR KEY:",
             size=7, bold=True, color=MID_GREY)
    for i, (sec, col) in enumerate(SECTOR_PAL.items()):
        kx = lx + 1.05 + i * 1.72
        add_rect(slide, kx, key_y+0.04, 0.13, 0.13, col)
        add_text(slide, kx+0.17, key_y, 1.5, 0.22, sec, size=7, color=MID_GREY)

    return slide

# ── SLIDE 2 ───────────────────────────────────────────────────────────────────
def build_slide2(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = 13.33, 7.5
    allocs = data["allocations"]
    grand  = sum(a["alloc"] for a in allocs)

    sector_totals = {}
    for a in allocs:
        sector_totals[a["sector"]] = sector_totals.get(a["sector"], 0) + a["alloc"]
    sorted_sectors = sorted(sector_totals.items(), key=lambda x: x[1], reverse=True)

    add_rect(slide, 0, 0, W, H, WHITE)

    # ── Market Context banner + text ─────────────────────────────────────────
    add_rect(slide, 0.3, 0.12, W-0.6, 0.30, GOLD_FILL)
    add_text(slide, 0.35, 0.12, 4, 0.30,
             "Market Context – " + data["reportDate"],
             size=10, bold=True, color=DARK_GREY)

    add_rect(slide, 0.3, 0.44, W-0.6, 0.70, LIGHT_GREY, BORDER, 0.3)
    add_text(slide, 0.42, 0.46, W-0.82, 0.66,
             data["marketContext"], size=10, color=DARK_GREY, wrap=True)

    # ── Left column: donut chart ──────────────────────────────────────────────
    chart_x, chart_y = 0.3, 1.30
    chart_w, chart_h = 6.1, 4.0

    # Sector Allocation gold banner
    add_rect(slide, chart_x, chart_y, chart_w, 0.32, GOLD_FILL)
    add_text(slide, chart_x, chart_y, chart_w, 0.32,
             "SECTOR ALLOCATION", size=10, bold=True,
             color=DARK_GREY, align=PP_ALIGN.CENTER)

    # Donut chart
    cd = ChartData()
    cd.categories = [s for s, v in sorted_sectors]
    cd.add_series("", [v for s, v in sorted_sectors])
    donut_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(chart_x), Inches(chart_y + 0.32),
        Inches(chart_w), Inches(chart_h - 0.32),
        cd
    )
    ch = donut_shape.chart
    ch.has_legend = True
    ch.legend.position = 2  # right
    ch.legend.include_in_layout = False
    # Show % labels on each segment
    dl = ch.series[0].data_labels
    dl.show_percentage = True
    dl.show_value      = False
    dl.number_format   = "0%"

    # Colour segments
    for i, (sec, _) in enumerate(sorted_sectors):
        pt = ch.series[0].points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = SECTOR_PAL.get(sec, RGBColor(0x88,0x88,0x88))

    # ── Right column: diversification blocks ─────────────────────────────────
    rx = chart_x + chart_w + 0.25
    rw = W - rx - 0.3
    ry = 1.30

    div_blocks = [
        ("GEOGRAPHIC DIVERSIFICATION",
         data.get("geographicDiversification", [])),
        ("SECTOR DIVERSIFICATION",
         data.get("sectorDiversification", [])),
        ("CURRENCY DIVERSIFICATION",
         data.get("currencyDiversification", [])),
    ]

    block_h   = (chart_h) / 3
    banner_h  = 0.30
    body_h    = block_h - banner_h - 0.04

    for i, (title, bullets) in enumerate(div_blocks):
        by = ry + i * (block_h + 0.04)
        add_rect(slide, rx, by, rw, banner_h, GOLD_FILL)
        add_text(slide, rx+0.06, by, rw-0.08, banner_h,
                 title, size=9, bold=True,
                 color=DARK_GREY, align=PP_ALIGN.CENTER)
        add_rect(slide, rx, by+banner_h, rw, body_h, WHITE, BORDER, 0.3)
        add_bullets(slide, rx+0.08, by+banner_h+0.05,
                    rw-0.14, body_h-0.08, bullets, size=8.5)

    # ── Investment Rationale ──────────────────────────────────────────────────
    rat_y = chart_y + chart_h + 0.10
    rat_h = H - rat_y - 0.12

    # Full width gold banner
    add_rect(slide, 0.3, rat_y, W-0.6, 0.28, GOLD_FILL)
    add_text(slide, 0.35, rat_y, W-0.6, 0.28,
             "INVESTMENT RATIONALE", size=10, bold=True, color=DARK_GREY)

    # Three columns
    col_w  = (W - 0.6) / 3
    col_h  = rat_h - 0.28 - 0.04
    cols   = [
        (data.get("rationaleCol1Title",""), data.get("rationaleCol1Body","")),
        (data.get("rationaleCol2Title",""), data.get("rationaleCol2Body","")),
        (data.get("rationaleCol3Title",""), data.get("rationaleCol3Body","")),
    ]
    for i, (title, body) in enumerate(cols):
        cx = 0.3 + i * col_w
        cy = rat_y + 0.28 + 0.04
        add_rect(slide, cx, cy, col_w-0.04, col_h, WHITE, BORDER, 0.3)
        add_text(slide, cx+0.1, cy+0.06, col_w-0.2, 0.28,
                 title, size=9, bold=True, color=DARK_GREY)
        add_text(slide, cx+0.1, cy+0.34, col_w-0.2, col_h-0.4,
                 body, size=8.5, color=MID_GREY, wrap=True)

    return slide

# ── Flask endpoints ────────────────────────────────────────────────────────────
@app.route("/generate", methods=["POST"])
def generate():
    data = request.get_json().get("reportData", {})
    prs  = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    build_slide1(prs, data)
    build_slide2(prs, data)
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return jsonify({"pptxBase64": base64.b64encode(buf.read()).decode()})

@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8080))
    app.run(host="0.0.0.0", port=port, debug=False)
